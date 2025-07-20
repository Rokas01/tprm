from pptx import Presentation
from pptx.util import Pt
from io import BytesIO

def apply_font(text_object):
    font = text_object.font
    font.name = 'Calibri'
    font.size = Pt(8)
    #font.bold = True
    #font.italic = None  #value to be inherited from theme


def add_slide(prs_object_input, title_input, requirement_text=True):

    slide_heading_from_call = title_input[0]
    requirement_text_from_call = title_input[1]
    summary_text_from_call = title_input[2]
    observation_text_from_call = title_input[3]
    risks_text_from_call = title_input[4]
    
    #Adding slide 1
    #==============
    slide1 = prs_object_input.slides.add_slide(prs_object_input.slide_layouts[1])

    heading_box, requirement_box, summary_box = slide1.placeholders

    title_para = slide1.shapes.title.text_frame.paragraphs[0]
    title_para.font.size = Pt(26)
    title_para.text = slide_heading_from_call

    if requirement_text == True:

        requirement_box.text = requirement_text_from_call

    summary_box.text = summary_text_from_call


    #Adding slide 2
    #==============
    slide2 = prs_object_input.slides.add_slide(prs_object_input.slide_layouts[2])

    title_para = slide2.shapes.title.text_frame.paragraphs[0]
    title_para.font.size = Pt(26)
    title_para.text = slide_heading_from_call

    heading_box, findings_box, risks_box = slide2.placeholders

    findings_box.text = observation_text_from_call
    
    risks_box.text = risks_text_from_call

def create_presentation_report_findings(title_text, subtitle_text, requirementsAndFindings, include_requirement_text=True, custom_template='template-NIS2.pptx'):

    prs = Presentation(custom_template)

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title_box, subtitle_box = slide.placeholders

    title_box.text = title_text
    subtitle_box.text = subtitle_text

    for requirement in requirementsAndFindings:

        add_slide(prs, requirement, requirement_text=include_requirement_text)
 
    prs.save("filename.pptx")

    binary_output = BytesIO()
    prs.save(binary_output) 

    return binary_output.getvalue()

#test_results = [["Article 21", "This is a requirement text This is a requirement text This", "this is a summary", "This is an observation", "This is a risk"], ["Article 22", "test text for the next article", "this is a summary", "not requested", "this is input 4"]]

#create_presentation_report_findings("NIS2 compliance", "assessment report", test_results)

